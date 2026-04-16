import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Define a Modern Dark Theme
BG_COLOR = RGBColor(18, 18, 20)           # Very Dark Grey
CARD_COLOR = RGBColor(30, 31, 35)         # Elevated Dark Grey
ACCENT_COLOR = RGBColor(94, 116, 255)     # Vibrant Indigo/Blue
ACCENT_LIGHT = RGBColor(140, 158, 255)    # Lighter Indigo
TEXT_MAIN = RGBColor(245, 245, 245)       # Off White
TEXT_MUTED = RGBColor(160, 160, 165)      # Muted Grey

# Define presentation
prs = Presentation()
# Set Aspect Ratio to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Clear default slides if any, actually we just append using blank layout (layout 6)
BLANK_LAYOUT = prs.slide_layouts[6]

def set_background(slide):
    """Sets a solid dark background for the slide with subtle decorative elements"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Top thin gradient-like accent bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT_COLOR
    top_bar.line.fill.background()

def add_title(slide, text):
    """Adds a styled title and an accent underline"""
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    p.font.name = "Inter" # Requires modern font, fallback is Arial
    
    # Accent underline
    line = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.3), Inches(1.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()

def populate_text_in_card(tf, content):
    """Formats bullet points nicely within a card's text frame"""
    tf.word_wrap = True
    lines = content.strip().split('\n')
    
    first = True
    for line in lines:
        if not line.strip():
            continue
            
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            
        p.font.name = "Inter"
        
        # Check if it's a main bullet or sub-bullet
        if line.startswith("•") or line.startswith("-"):
            p.text = line[1:].strip()
            # Bullet point settings
            p.level = 1 if line.startswith("-") else 0
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_MAIN
            p.space_after = Pt(12)
            p.space_before = Pt(6)
        else:
            # Regular text (might be a sub-heading inside card)
            p.text = line
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = ACCENT_LIGHT
            p.space_after = Pt(8)
            p.space_before = Pt(8)

slides_data = [
    {
        "type": "title",
        "title": "Smart Air Quality\nMonitoring & Alert System",
        "subtitle": "Course: BCSE301L Software Engineering\nFaculty: S. Margret Anouncia | VIT University\n\nTeam 14: Yash Raj • Laksham Singh • Himansh Mangal • Abhay Mittal • Fabrice Freejo"
    },
    {
        "type": "content",
        "title": "Project Overview",
        "content": "Introduction\n• A next-generation software solution designed for real-time environmental monitoring.\nSocietal Problem\n• Rising health risks (respiratory/cardiovascular) due to lack of real-time, actionable pollution data in urban zones.\nTarget Users\n• Citizens, Environmental Agencies, Municipal Authorities, and Healthcare Organizations."
    },
    {
        "type": "content",
        "title": "Project Objectives",
        "content": "Performance-based Goals\n• To collect air quality data from multiple locations with high accuracy.\n• To monitor pollution levels continuously with minimal latency.\n• To display real-time AQI data with fast response time.\n• To generate alerts instantly when pollution exceeds safe limits.\n• To store historical data efficiently for reporting."
    },
    {
        "type": "split",
        "title": "Project Planning (Process WBS)",
        "content": "Phases Framework\n• Planning & Design\n• Installation & Setup\n• Data Collection & Analysis\n• Alert & Response",
        "image_placeholder": "PROCESS-BASED WBS DIAGRAM\n\nDrop Image Here\n(From Tool implementation Of Task-1.pdf, Page 1)"
    },
    {
        "type": "split",
        "title": "Project Planning (Product WBS)",
        "content": "Core Deliverables\n• Comprehensive Documentation\n• Sensor Networks & Nodes\n• Data Processing Module\n• Alert & Notification Module",
        "image_placeholder": "PRODUCT-BASED WBS DIAGRAM\n\nDrop Image Here\n(From Tool implementation Of Task-1.pdf, Page 3)"
    },
    {
        "type": "split",
        "title": "Gantt Chart: Process-Based",
        "content": "Strategic Timeline\n• February 2026 – April 2026\n• Encompasses project initiation through database design and dashboard development.",
        "image_placeholder": "PROCESS-BASED GANTT CHART\n\nDrop Image Here\n(From Tool implementation Of Task-1.pdf, Page 6)"
    },
    {
        "type": "split",
        "title": "Gantt Chart: Product-Based",
        "content": "Delivery Timeline\n• February 2026 – April 2026\n• Covers system specs, UI design, sensor inputs, and final testing evaluations.",
        "image_placeholder": "PRODUCT-BASED GANTT CHART\n\nDrop Image Here\n(From Tool implementation Of Task-1.pdf, Page 8)"
    },
    {
        "type": "content",
        "title": "Requirement Elicitation",
        "content": "Methodology\n• Active Observation across residential, educational, and traffic-heavy industrial environments.\nKey Findings\n• Citizens face heightened health risks largely stemming from a lack of immediate alert systems.\n• The data is either complex or completely unavailable to the general public."
    },
    {
        "type": "content",
        "title": "Viewpoint-Oriented Analysis",
        "content": "Citizen Viewpoint\n• Needs simplified real-time AQI data and instant safety alerts directly to mobile devices.\nAuthority Viewpoint\n• Requires environmental trend tracking, city-wide analysis, and actionable policy reporting.\nSystem Administrator\n• Needs a reliable dashboard for sensor maintenance, fault tracking, and node management."
    },
    {
        "type": "content",
        "title": "System Design Overview",
        "content": "Strategic Approach\n• Focused on high modularity, uncompromised real-time data collection, and efficient processing of raw parameters into Actionable AQI.\nToolset\n• Engineered and modeled using industry-standard tools: Modelio 5.4 & Draw.io."
    },
    {
        "type": "split",
        "title": "Detailed Object Design",
        "content": "Core Entity Classes\n• Sensor Module\n• Data Processing Unit\n• Alert & Notification System\n• Central Database Ecosystem\n• Client Application (User)",
        "image_placeholder": "CLASS DIAGRAM\n\nDrop Image Here\n(From Task-3 Tool Work.pdf, Page 1)"
    },
    {
        "type": "split",
        "title": "System Interaction",
        "content": "Use Case Actors\n• Citizen (User)\n• Environmental Authority\n• System Administrator\nCore Functionalities\n• Monitor Real-Time AQI\n• Receive Instant Notifications\n• Evaluate System Performance",
        "image_placeholder": "USE CASE DIAGRAM\n\nDrop Image Here"
    },
    {
        "type": "split",
        "title": "Architectural Blueprint",
        "content": "Pattern\n• Layered Architecture ensuring strict isolation of concerns.\nActive Layers\n• Presentation Layer (UI)\n• Business Logic Layer\n• Persistence Layer\n• Core Database Layer",
        "image_placeholder": "LAYERED ARCHITECTURE DIAGRAM\n\nDrop Image Here\n(From Task-3 Tool Work.pdf, Page 3)"
    },
    {
        "type": "split",
        "title": "Cohesion & Coupling Analysis",
        "content": "Metrics Evaluated\n• High Functional Cohesion: Guarantees each component executes a highly focused specific task.\n• Low Data Coupling: Modules communicate via optimized, minimalist data payloads.\nOutcome\n• A robust, fault-tolerant system ready for integration.",
        "image_placeholder": "COHESION & COUPLING DIAGRAM\n\nDrop Image Here\n(From Task-3 Tool Work.pdf, Page 4)"
    },
    {
        "type": "content",
        "title": "Testing Strategy",
        "content": "Methodologies Deployed\n• Boundary Value Analysis (Evaluating edge thresholds)\n• Equivalence Partitioning (Testing standardized input ranges)\nExecutive Summary\n• Successfully validated TC-01 to TC-06.\n• Achieved an uncompromising 100% Pass Rate across varying scenarios."
    },
    {
        "type": "content",
        "title": "Conclusion & Future Scope",
        "content": "Retrospective Summary\n• Successfully architected a reliable, modular framework delivering actionable AQI monitoring and instantaneous safety alerts.\nFuture Roadmap\n• Implementation of AI-based predictive forecasting.\n• Deep integration into larger Smart City frameworks.\n• Hyper-personalized health & route safety recommendations."
    }
]

# Generate Slides
for data in slides_data:
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_background(slide)
    
    if data["type"] == "title":
        # Add decorative circle
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(1), Inches(5), Inches(5))
        oval.fill.solid()
        oval.fill.fore_color.rgb = ACCENT_COLOR
        oval.line.fill.background()
        
        # Transparent overlay for styling
        overlay = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(1.5), Inches(5), Inches(5))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = BG_COLOR
        overlay.line.fill.background()

        # Main Title Card
        title_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.8), Inches(10), Inches(4))
        title_card.fill.solid()
        title_card.fill.fore_color.rgb = CARD_COLOR
        title_card.line.fill.background()
        
        tf = title_card.text_frame
        tf.margin_top = Inches(0.5)
        tf.margin_left = Inches(0.5)
        
        p = tf.paragraphs[0]
        p.text = data["title"]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        p.font.name = "Inter"
        
        # Add a gap
        p2 = tf.add_paragraph()
        p2.space_before = Pt(20)
        
        for line in data["subtitle"].split('\n'):
            p_sub = tf.add_paragraph()
            p_sub.text = line
            p_sub.font.size = Pt(16)
            p_sub.font.color.rgb = TEXT_MUTED
            p_sub.font.name = "Inter"

    elif data["type"] == "content":
        add_title(slide, data["title"])
        
        # Wide content card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLOR
        card.line.fill.background()
        card.text_frame.margin_left = Inches(0.5)
        card.text_frame.margin_top = Inches(0.5)
        card.text_frame.margin_right = Inches(0.5)
        
        populate_text_in_card(card.text_frame, data["content"])

    elif data["type"] == "split":
        add_title(slide, data["title"])
        
        # Left Text Card
        card_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
        card_left.fill.solid()
        card_left.fill.fore_color.rgb = CARD_COLOR
        card_left.line.fill.background()
        card_left.text_frame.margin_left = Inches(0.4)
        card_left.text_frame.margin_top = Inches(0.4)
        card_left.text_frame.margin_right = Inches(0.4)
        
        populate_text_in_card(card_left.text_frame, data["content"])
        
        # Right Image Placeholder Card
        card_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.8), Inches(6.0), Inches(4.8))
        card_right.fill.solid()
        card_right.fill.fore_color.rgb = RGBColor(24, 25, 28) # Slightly darker to indicate drop area
        card_right.line.color.rgb = ACCENT_COLOR # Give it a nice border
        card_right.line.width = Pt(1.5)
        
        tf2 = card_right.text_frame
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, line in enumerate(data["image_placeholder"].split('\n')):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.text = line
            p.font.size = Pt(20 if i == 0 else 14)
            p.font.bold = (i == 0)
            p.font.color.rgb = ACCENT_LIGHT if i == 0 else TEXT_MUTED

prs.save("Smart_Air_Quality_System_Premium.pptx")
print("Presentation successfully created: Smart_Air_Quality_System_Premium.pptx")
